# import streamlit as st
# from vanna.qianwen import QianWenAI_Chat
# from vanna.chromadb import ChromaDB_VectorStore

# 系统参数

import os
from dotenv import load_dotenv, find_dotenv
# import streamlit as st

# vanna
from vanna_main.qianwen import QianWenAI_Chat
from vanna_main.chromadb import ChromaDB_VectorStore


#我的共通
from common.my_config import MyConfig
config = MyConfig()



# @st.cache_resource(ttl=7860)
def setup_vanna():

    _ = load_dotenv(find_dotenv())
    DASHSCOPE_API_KEY = os.environ["DASHSCOPE_API_KEY"]

    class MyVanna(ChromaDB_VectorStore, QianWenAI_Chat):
        def __init__(self, config=None):
            ChromaDB_VectorStore.__init__(self, config=config)
            QianWenAI_Chat.__init__(self, config={'api_key': DASHSCOPE_API_KEY, 'model': 'qwen-plus'})

    vn = MyVanna()

    vn.connect_to_sqlite(config.work_db_path)

    return vn


# @st.cache_data(show_spinner="SQL生成中 ...")
def generate_sql_cached(question: str):
    vn = setup_vanna()
    return vn.generate_sql(question=question, allow_llm_to_see_data=True)


# @st.cache_data(show_spinner="SQL検証 ...")
def is_sql_valid_cached(sql: str):
    vn = setup_vanna()
    print("12" * 50)
    print("sql", sql)
    print("12" * 50)
    return vn.is_sql_valid(sql=sql)


# @st.cache_data(show_spinner="SQL実行 ...")
def run_sql_cached(sql: str):
    vn = setup_vanna()
    return vn.run_sql(sql=sql)


# @st.cache_data(show_spinner="チャット生成できるかどうか ...")
def should_generate_chart_cached(question, sql, df):
    vn = setup_vanna()
    return vn.should_generate_chart(df=df)


# @st.cache_data(show_spinner="チャット生成用のPythonコード生成 ...")
def generate_plotly_code_cached(question, sql, df):
    vn = setup_vanna()
    code = vn.generate_plotly_code(question=question, sql=sql, df=df)
    return code


# @st.cache_data(show_spinner="チャット生成 ...")
def generate_plot_cached(code, df):
    vn = setup_vanna()
    return vn.get_plotly_figure(plotly_code=code, df=df)


# @st.cache_data(show_spinner="要約生成 ...")
def generate_summary_cached(question, df):
    vn = setup_vanna()
    return vn.generate_summary(question=question, df=df)


# @st.cache_data(show_spinner="Generating sample questions ...")
def generate_questions_cached():
    vn = setup_vanna()
    return vn.generate_questions()


# @st.cache_data(show_spinner="Generating followup questions ...")
def generate_followup_cached(question, sql, df):
    vn = setup_vanna()
    return vn.generate_followup_questions(question=question, sql=sql, df=df)


# @st.cache_data(show_spinner="トレーニングデータの追加 ...")
def add_training_data_cached(question, sql):
    vn = setup_vanna()
    return vn.train(question=question, sql=sql)



# @st.cache_data(show_spinner="PPTファイルの出力 ...")
def save_ppt_cached(fig,summary):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from io import BytesIO
    from PIL import Image
    if fig is None:
        return
    if summary is None:
        return
    # Load the existing PowerPoint presentation
    ppt_path = "../ppt/Template_Oracle.pptx"  # Replace with your file path
    prs = Presentation(ppt_path)

    # Add a new slide (using a blank slide layout)
    slide_layout = prs.slide_layouts[10]  # 5 is usually a blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add text to the slide
    text_box = slide.shapes.add_textbox(Inches(3), Inches(6), Inches(8), Inches(1))

    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.text = summary
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)  # Set text color (blue)
    p.alignment = PP_ALIGN.LEFT  # Center-align the text

    # Load the image file into a BytesIO object
    image_path = "../pic/chart.png"  # Replace with the path to your image
    with open(image_path, "rb") as img_file:
        img_bytes = BytesIO(img_file.read())

    # Add the image to the slide
    slide.shapes.add_picture(img_bytes, Inches(3), Inches(1), width=Inches(7), height=Inches(5))

    # Save the modified PowerPoint
    output_path = "../ppt/ouput_ppt.pptx"
    prs.save(output_path)

    return

